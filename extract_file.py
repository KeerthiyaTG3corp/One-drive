import tempfile
import os

# PDF
import pdfplumber

# DOCX
from docx import Document

# Images
from PIL import Image
import pytesseract

# TXT
def extract_text_from_txt(b):
    return b.decode("utf-8", errors="ignore").strip()

# PDF
def extract_text_from_pdf(b):
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tf:
        tf.write(b)
        tf.flush()

    text = ""
    with pdfplumber.open(tf.name) as pdf:
        for page in pdf.pages:
            text += page.extract_text() or ""
    return text.strip()

# DOCX
def extract_text_from_docx(b):
    with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tf:
        tf.write(b)
        tf.flush()

    doc = Document(tf.name)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

# Image OCR
def extract_text_from_image(b):
    with tempfile.NamedTemporaryFile(delete=False) as tf:
        tf.write(b)
        tf.flush()

    img = Image.open(tf.name)
    return pytesseract.image_to_string(img).strip()

# XLSX (optional)
def extract_text_from_xlsx(b):
    import openpyxl
    with tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx") as tf:
        tf.write(b)
        tf.flush()

    wb = openpyxl.load_workbook(tf.name)
    text = ""
    for sheet in wb:
        for row in sheet.iter_rows(values_only=True):
            text += " ".join(str(c) for c in row if c) + "\n"
    return text.strip()

# PPTX (optional)
def extract_text_from_pptx(b):
    from pptx import Presentation
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tf:
        tf.write(b)
        tf.flush()

    prs = Presentation(tf.name)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text.strip()

# Universal dispatcher
def extract_text(file_bytes, filename):
    ext = os.path.splitext(filename.lower())[1]

    if ext == ".pdf":
        return extract_text_from_pdf(file_bytes)
    elif ext == ".docx":
        return extract_text_from_docx(file_bytes)
    elif ext in [".png", ".jpg", ".jpeg"]:
        return extract_text_from_image(file_bytes)
    elif ext == ".txt":
        return extract_text_from_txt(file_bytes)
    elif ext == ".xlsx":
        return extract_text_from_xlsx(file_bytes)
    elif ext == ".pptx":
        return extract_text_from_pptx(file_bytes)
    else:
        return ""  # fail-safe: ignore unsupported files
